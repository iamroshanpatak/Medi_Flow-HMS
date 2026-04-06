from __future__ import annotations

from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _set_title(slide, text: str) -> None:
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True


def _add_subtitle(slide, text: str) -> None:
    placeholder = slide.placeholders[1]
    placeholder.text = text
    p = placeholder.text_frame.paragraphs[0]
    p.font.size = Pt(18)


def _add_bullets(slide, title: str, bullets: list[str]) -> None:
    layout = slide.part.slide_layout
    _ = layout
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)


def _add_two_column(slide, title: str, left_title: str, left: list[str], right_title: str, right: list[str]) -> None:
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(34)
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True

    # Remove the default content placeholder if present
    for shape in list(slide.shapes):
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            sp = shape._element
            sp.getparent().remove(sp)

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(4.5), Inches(5.0))
    right_box = slide.shapes.add_textbox(Inches(5.1), Inches(1.7), Inches(4.7), Inches(5.0))

    def fill(box, header: str, items: list[str]) -> None:
        tf = box.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = header
        p0.font.size = Pt(24)
        p0.font.bold = True
        p0.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)  # slate-ish

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)

    fill(left_box, left_title, left)
    fill(right_box, right_title, right)


def build_presentation(output_path: Path) -> None:
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_title(slide, "MediFlow")
    _add_subtitle(
        slide,
        f"Hospital Management System\nPresentation | {date.today().isoformat()}",
    )

    # Agenda
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Agenda",
        [
            "Problem & objectives",
            "System overview (roles & modules)",
            "Architecture & tech stack",
            "Key workflows (appointments + queue)",
            "Security & database",
            "Testing & conclusion",
        ],
    )

    # Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Problem Statement",
        [
            "Manual hospital processes cause delays and confusion",
            "Patients cannot easily track token/queue status",
            "Doctors need a live view of who is next",
            "Medical info can be scattered or hard to access",
        ],
    )

    # Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Project Objectives",
        [
            "Online appointment booking with time slots",
            "Real-time queue tracking (token, position, wait time)",
            "Role-based dashboards: Patient / Doctor / Admin",
            "Medical records storage and viewing",
            "Secure authentication and authorization",
        ],
    )

    # Scope
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Scope (What the system covers)",
        [
            "Patients: register/login, book appointments, check queue status",
            "Doctors: manage queue, call next, complete consultation",
            "Admins/Staff: monitor system and support operations",
            "Database: users, appointments, queue, medical records",
        ],
    )

    # Roles & features
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    _add_two_column(
        slide,
        "User Roles & Main Features",
        "Patient",
        [
            "Book appointment",
            "Check-in and receive token",
            "Track queue position",
            "View medical records",
        ],
        "Doctor / Admin",
        [
            "See live queue",
            "Call next patient",
            "Complete consultation",
            "Manage availability (doctor)",
        ],
    )

    # Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "System Architecture (High Level)",
        [
            "Frontend: Next.js (React + TypeScript) UI",
            "Backend: Node.js + Express REST API",
            "Real-time: Socket.IO for queue updates",
            "Database: MongoDB with Mongoose models",
        ],
    )

    # Tech stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Tech Stack",
        [
            "Frontend: Next.js, React, TypeScript, Tailwind CSS",
            "Backend: Express.js, JWT, bcryptjs",
            "Database: MongoDB, Mongoose",
            "Real-time: Socket.IO",
            "Tools: npm, nodemon, ESLint",
        ],
    )

    # Workflow: appointment booking
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Workflow 1: Appointment Booking",
        [
            "Patient selects doctor",
            "Select date and available time slot",
            "Submit reason and create appointment",
            "System prevents conflicting bookings",
        ],
    )

    # Workflow: queue management
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Workflow 2: Real-Time Queue",
        [
            "Patient checks-in on appointment day to get token",
            "Doctor sees queue updates instantly",
            "Doctor calls next → patient gets notification",
            "Complete consultation → statuses update",
        ],
    )

    # Security
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Security",
        [
            "JWT-based login",
            "Protected routes and API authorization by role",
            "Passwords hashed using bcryptjs",
            "Patients can only access their own records",
        ],
    )

    # Database (simple)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Database (Main Collections)",
        [
            "Users: patient/doctor/admin/staff profiles",
            "Appointments: doctor + patient + date/time + status",
            "Queue: token number, position, wait time, status",
            "MedicalRecords: vitals, diagnosis, prescriptions, notes",
        ],
    )

    # Testing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Testing (What we verified)",
        [
            "Login/registration works and routes are protected",
            "Booking creates appointments and avoids conflicts",
            "Queue updates in real time (doctor + patient)",
            "Medical records CRUD respects permissions",
        ],
    )

    # Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _add_bullets(
        slide,
        "Conclusion & Future Work",
        [
            "MediFlow improves scheduling and reduces waiting confusion",
            "Real-time queue provides better patient experience",
            "Future: stronger notifications, reports, audit logs, CI tests",
        ],
    )

    # Thank you
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_title(slide, "Thank You")
    box = slide.shapes.add_textbox(Inches(1.2), Inches(2.4), Inches(8.0), Inches(2.0))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Questions?"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


if __name__ == "__main__":
    root = Path(__file__).resolve().parents[1]
    out = root / "docs" / "MediFlow_Presentation.pptx"
    build_presentation(out)
    print(f"Created: {out}")
